import logging
import re
import zipfile
import xml.etree.ElementTree as ET
from io import BytesIO
from pathlib import Path

import fitz
import openpyxl
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Emu

from app.config import settings
from app.converters.result import ConversionResult
from app.errors import ConversionError

# 300 DPI is the standard print-quality baseline. PNG remains lossless; JPEG
# uses a high quality setting to avoid the visible artifacts produced by the
# encoder default while keeping downloads practical.
RASTER_DPI = 300
JPEG_QUALITY = 95
_ILLEGAL_XLSX_CHARS = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")

# Below this average chars-per-page we treat the PDF as scanned (image-only).
_SCANNED_THRESHOLD = 30


def _open_pdf(input_path: Path) -> fitz.Document:
    try:
        doc = fitz.open(input_path)
    except Exception as exc:
        raise ConversionError(
            422, "unreadable_file", "The PDF could not be read. It may be corrupted."
        ) from exc

    if not doc.is_pdf or doc.page_count == 0:
        doc.close()
        raise ConversionError(422, "unreadable_file", "The uploaded file is not a valid PDF.")

    return doc


def _rasterize(input_path: Path, out_dir: Path, ext: str, media_type: str) -> ConversionResult:
    doc = _open_pdf(input_path)
    try:
        paths = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=RASTER_DPI)
            img_path = out_dir / f"page-{i + 1}{ext}"
            if ext.lower() in {".jpg", ".jpeg"}:
                pix.save(img_path, jpg_quality=JPEG_QUALITY)
            else:
                pix.save(img_path)
            paths.append(img_path)
    finally:
        doc.close()

    if len(paths) == 1:
        return ConversionResult(paths[0], media_type)

    zip_path = out_dir / "output.zip"
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for p in paths:
            zf.write(p, p.name)
    return ConversionResult(zip_path, "application/zip")


def convert_pdf_to_jpg(input_path: Path, out_dir: Path) -> ConversionResult:
    return _rasterize(input_path, out_dir, ".jpg", "image/jpeg")


def convert_pdf_to_png(input_path: Path, out_dir: Path) -> ConversionResult:
    return _rasterize(input_path, out_dir, ".png", "image/png")


def convert_pdf_to_text(input_path: Path, out_dir: Path) -> ConversionResult:
    doc = _open_pdf(input_path)
    try:
        text = "\n\n".join(page.get_text() for page in doc)
    finally:
        doc.close()

    output_path = out_dir / "output.txt"
    output_path.write_text(text, encoding="utf-8")
    return ConversionResult(output_path, "text/plain")


def _avg_chars_per_page(doc: fitz.Document) -> float:
    """Return average extractable character count per page."""
    if doc.page_count == 0:
        return 0.0
    return sum(len(page.get_text().strip()) for page in doc) / doc.page_count


def convert_pdf_to_docx(
    input_path: Path,
    out_dir: Path,
    mode: str = "keep-layout",
) -> ConversionResult:
    # Validate and check for scanned (image-only) content before attempting conversion.
    doc = _open_pdf(input_path)
    avg_chars = _avg_chars_per_page(doc)

    if avg_chars < _SCANNED_THRESHOLD:
        doc.close()
        raise ConversionError(
            422,
            "scanned_pdf",
            "This PDF appears to be scanned — it contains no selectable text. "
            "Please run it through an OCR tool first, then convert the resulting PDF.",
        )

    output_path = out_dir / f"{input_path.stem}.docx"
    try:
        from app.converters.pdf_word.docx_builder import (
            build_docx,
            build_keep_layout_docx,
        )

        if mode == "keep-layout":
            build_keep_layout_docx(doc, output_path)
        elif mode == "editable":
            from app.converters.pdf_word.layout_analyzer import analyze_document

            layouts = analyze_document(doc)
            build_docx(layouts, output_path)
        else:
            raise ConversionError(
                400,
                "invalid_conversion_mode",
                "PDF to Word mode must be 'keep-layout' or 'editable'.",
            )
    except ConversionError:
        raise
    except Exception as exc:
        raise ConversionError(
            422,
            "unreadable_file",
            "The PDF could not be converted. It may be corrupted or use an unsupported format.",
        ) from exc
    finally:
        doc.close()

    if not output_path.is_file():
        raise ConversionError(500, "conversion_failed", "Conversion did not produce an output file.")

    return ConversionResult(
        output_path,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    )


def convert_pdf_to_xlsx(input_path: Path, out_dir: Path) -> ConversionResult:
    doc = _open_pdf(input_path)
    thumbnail = doc[0].get_pixmap(dpi=72, alpha=False).tobytes("jpeg")
    try:
        if doc.page_count > settings.max_pdf_pages:
            raise ConversionError(
                422,
                "too_many_pages",
                f"PDF exceeds the {settings.max_pdf_pages}-page conversion limit.",
            )
        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # remove default empty sheet
        for i, page in enumerate(doc):
            ws = wb.create_sheet(title=f"Page {i + 1}")
            tables = page.find_tables().tables
            next_row = 1
            if tables:
                for table_index, table in enumerate(tables):
                    rows = table.extract()
                    for row_index, row in enumerate(rows):
                        for col_index, value in enumerate(row, start=1):
                            text = _safe_excel_text(value)
                            cell = ws.cell(next_row + row_index, col_index, text)
                            cell.alignment = Alignment(vertical="top", wrap_text=True)
                            if row_index == 0:
                                cell.font = Font(bold=True)
                                cell.fill = PatternFill("solid", fgColor="D9EAF7")
                    next_row += len(rows) + 2
            else:
                # Useful fallback for PDFs without explicit table geometry.
                for line in page.get_text("text").splitlines():
                    if line.strip():
                        ws.cell(next_row, 1, _safe_excel_text(line))
                        next_row += 1

            ws.freeze_panes = "A2"
            for column in range(1, min(ws.max_column, 100) + 1):
                max_len = max(
                    (len(str(ws.cell(row, column).value or "")) for row in range(1, ws.max_row + 1)),
                    default=0,
                )
                ws.column_dimensions[get_column_letter(column)].width = min(max(max_len + 2, 10), 60)
    finally:
        doc.close()

    output_path = out_dir / f"{input_path.stem}.xlsx"
    wb.save(output_path)
    _inject_ooxml_thumbnail(output_path, thumbnail)
    return ConversionResult(
        output_path,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    )


def _safe_excel_text(value) -> str:
    text = _ILLEGAL_XLSX_CHARS.sub("", "" if value is None else str(value))
    # Prevent formulas embedded in untrusted PDF text from executing in Excel.
    if text.lstrip().startswith(("=", "+", "-", "@")):
        return "'" + text
    return text


def convert_pdf_to_pptx(input_path: Path, out_dir: Path) -> ConversionResult:
    doc = _open_pdf(input_path)
    thumbnail = doc[0].get_pixmap(dpi=72, alpha=False).tobytes("jpeg")
    try:
        if doc.page_count > settings.max_pdf_pages:
            raise ConversionError(
                422,
                "too_many_pages",
                f"PDF exceeds the {settings.max_pdf_pages}-page conversion limit.",
            )

        presentation = Presentation()
        # Remove the template slide and match the PDF's first page dimensions.
        for slide_id in list(presentation.slides._sldIdLst):
            presentation.part.drop_rel(slide_id.rId)
            presentation.slides._sldIdLst.remove(slide_id)
        first = doc[0].rect
        presentation.slide_width = Emu(round(first.width * 12700))
        presentation.slide_height = Emu(round(first.height * 12700))
        blank_layout = presentation.slide_layouts[6]

        for page in doc:
            slide = presentation.slides.add_slide(blank_layout)
            pix = page.get_pixmap(dpi=RASTER_DPI, alpha=False)
            slide.shapes.add_picture(
                BytesIO(pix.tobytes("png")),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )
    finally:
        doc.close()

    output_path = out_dir / f"{input_path.stem}.pptx"
    presentation.save(output_path)
    _inject_ooxml_thumbnail(output_path, thumbnail)
    return ConversionResult(
        output_path,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


def _inject_ooxml_thumbnail(package_path: Path, jpeg_data: bytes) -> None:
    """Add the standard OOXML package thumbnail used by Explorer/Office."""
    content_types_ns = "http://schemas.openxmlformats.org/package/2006/content-types"
    rels_ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    thumbnail_type = (
        "http://schemas.openxmlformats.org/package/2006/relationships/metadata/thumbnail"
    )
    ET.register_namespace("", content_types_ns)
    ET.register_namespace("", rels_ns)

    temp_path = package_path.with_suffix(package_path.suffix + ".tmp")
    with zipfile.ZipFile(package_path, "r") as source, zipfile.ZipFile(
        temp_path, "w", zipfile.ZIP_DEFLATED
    ) as target:
        content_types = ET.fromstring(source.read("[Content_Types].xml"))
        if not any(
            child.get("Extension", "").lower() in {"jpg", "jpeg"}
            for child in content_types
        ):
            ET.SubElement(
                content_types,
                f"{{{content_types_ns}}}Default",
                Extension="jpeg",
                ContentType="image/jpeg",
            )

        relationships = ET.fromstring(source.read("_rels/.rels"))
        if not any(child.get("Type") == thumbnail_type for child in relationships):
            used_ids = {child.get("Id") for child in relationships}
            index = 1
            while f"rIdThumb{index}" in used_ids:
                index += 1
            ET.SubElement(
                relationships,
                f"{{{rels_ns}}}Relationship",
                Id=f"rIdThumb{index}",
                Type=thumbnail_type,
                Target="docProps/thumbnail.jpeg",
            )

        replaced = {"[Content_Types].xml", "_rels/.rels", "docProps/thumbnail.jpeg"}
        for item in source.infolist():
            if item.filename not in replaced:
                target.writestr(item, source.read(item.filename))
        target.writestr(
            "[Content_Types].xml",
            ET.tostring(content_types, encoding="utf-8", xml_declaration=True),
        )
        target.writestr(
            "_rels/.rels",
            ET.tostring(relationships, encoding="utf-8", xml_declaration=True),
        )
        target.writestr("docProps/thumbnail.jpeg", jpeg_data)
    temp_path.replace(package_path)
