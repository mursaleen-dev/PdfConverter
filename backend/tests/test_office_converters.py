from pathlib import Path
import zipfile

import fitz
import openpyxl
import pytest
from pptx import Presentation
from pptx.util import Inches

from app.config import SOFFICE_PATH
from app.converters.office_converter import convert_office_to_pdf
from app.converters.pdf_converter import (
    _safe_excel_text,
    convert_pdf_to_pptx,
    convert_pdf_to_xlsx,
)


def _table_pdf(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=400, height=300)
    xs, ys = [40, 170, 300], [50, 80, 110]
    for x in xs:
        page.draw_line((x, ys[0]), (x, ys[-1]))
    for y in ys:
        page.draw_line((xs[0], y), (xs[-1], y))
    for x, y, text in (
        (50, 70, "Name"),
        (180, 70, "Amount"),
        (50, 100, "Widget"),
        (180, 100, "42"),
    ):
        page.insert_text((x, y), text, fontsize=10)
    doc.save(path)
    doc.close()


def test_pdf_to_excel_extracts_table(tmp_path):
    source = tmp_path / "table.pdf"
    _table_pdf(source)
    result = convert_pdf_to_xlsx(source, tmp_path)
    workbook = openpyxl.load_workbook(result.path)
    sheet = workbook["Page 1"]
    assert sheet["A1"].value == "Name"
    assert sheet["B2"].value == "42"
    assert sheet["A1"].font.bold
    with zipfile.ZipFile(result.path) as package:
        assert "docProps/thumbnail.jpeg" in package.namelist()


def test_pdf_to_excel_blocks_formula_injection():
    assert _safe_excel_text("=HYPERLINK(\"bad\")").startswith("'")
    assert _safe_excel_text("normal") == "normal"


def test_pdf_to_powerpoint_preserves_page_as_slide(tmp_path):
    source = tmp_path / "slides.pdf"
    _table_pdf(source)
    result = convert_pdf_to_pptx(source, tmp_path)
    presentation = Presentation(result.path)
    assert len(presentation.slides) == 1
    assert len(presentation.slides[0].shapes) == 1
    assert presentation.slides[0].shapes[0].shape_type == 13  # picture
    with zipfile.ZipFile(result.path) as package:
        assert "docProps/thumbnail.jpeg" in package.namelist()


@pytest.mark.skipif(not SOFFICE_PATH, reason="LibreOffice unavailable")
@pytest.mark.parametrize("kind", ["excel", "powerpoint"])
def test_office_to_pdf_end_to_end(tmp_path, kind):
    if kind == "excel":
        source = tmp_path / "sheet.xlsx"
        workbook = openpyxl.Workbook()
        sheet = workbook.active
        sheet.append(["Product", "Amount"])
        sheet.append(["Widget", 42])
        workbook.save(source)
    else:
        source = tmp_path / "slides.pptx"
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
        box.text = "Production conversion"
        presentation.save(source)

    result = convert_office_to_pdf(source, tmp_path)
    with fitz.open(result.path) as pdf:
        assert pdf.page_count >= 1
        assert pdf.is_pdf
